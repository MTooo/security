#!/usr/bin/env python3
"""生成 SM2 安全数据交换 SDK 技术分享 PPT"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# ── 颜色方案 ──────────────────────────────────────────
C_DARK   = RGBColor(0x1A, 0x1F, 0x36)   # 深蓝黑背景
C_PRIMARY = RGBColor(0x2D, 0x5B, 0xD7)  # 主色蓝
C_ACCENT  = RGBColor(0x00, 0xC9, 0xA7)  # 强调绿/青
C_WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
C_LIGHT   = RGBColor(0xCC, 0xD5, 0xE8)  # 浅灰蓝
C_WARN    = RGBColor(0xFF, 0x9F, 0x43)  # 橙色
C_RED     = RGBColor(0xFF, 0x6B, 0x6B)  # 红色
C_CODE_BG = RGBColor(0x0D, 0x11, 0x1E)  # 代码块背景
C_GRAY    = RGBColor(0x8A, 0x94, 0xA6)
C_DARK2   = RGBColor(0x25, 0x2D, 0x4A)
C_TABLE_H = RGBColor(0x2D, 0x5B, 0xD7)
C_TABLE_R = RGBColor(0x25, 0x2D, 0x4A)
C_BORDER  = RGBColor(0x3A, 0x4A, 0x6A)

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)

W = prs.slide_width
H = prs.slide_height

# ── 工具函数 ──────────────────────────────────────────

def add_bg(slide, color=C_DARK):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_rect(slide, left, top, width, height, fill_color=None, border_color=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.line.fill.background()
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    return shape

def add_text_box(slide, left, top, width, height, text="", font_size=18,
                 color=C_WHITE, bold=False, alignment=PP_ALIGN.LEFT, font_name="微软雅黑"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return tf

def add_para(tf, text, font_size=16, color=C_WHITE, bold=False, alignment=PP_ALIGN.LEFT,
             space_before=Pt(4), space_after=Pt(2), font_name="微软雅黑"):
    p = tf.add_paragraph()
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    p.space_before = space_before
    p.space_after = space_after
    return p

def add_code_block(slide, left, top, width, height, code_text, font_size=11):
    """添加代码块"""
    shape = add_rect(slide, left, top, width, height, C_CODE_BG)
    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = Pt(12)
    tf.margin_top = Pt(10)
    tf.margin_right = Pt(12)
    tf.margin_bottom = Pt(10)
    lines = code_text.strip().split('\n')
    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(font_size)
        p.font.name = "Consolas"
        p.font.color.rgb = C_ACCENT
        p.space_before = Pt(1)
        p.space_after = Pt(1)
    return shape

def add_section_header(slide, number, title):
    """章节标题栏"""
    add_rect(slide, Inches(0), Inches(0), W, Inches(1.1), C_PRIMARY)
    add_text_box(slide, Inches(0.6), Inches(0.2), Inches(12), Inches(0.7),
                 f"{'0' if number < 10 else ''}{number}  {title}",
                 font_size=32, color=C_WHITE, bold=True)

def add_page_number(slide, num, total=""):
    txt = f"{num}" if not total else f"{num} / {total}"
    add_text_box(slide, Inches(12.2), Inches(7.0), Inches(1), Inches(0.4),
                 txt, font_size=10, color=C_GRAY, alignment=PP_ALIGN.RIGHT)

def add_bottom_line(slide):
    add_rect(slide, Inches(0.5), Inches(7.1), Inches(12.3), Pt(2), C_PRIMARY)

def add_table(slide, left, top, col_widths, headers, rows, font_size=11):
    """添加格式化表格"""
    n_rows = len(rows) + 1
    n_cols = len(headers)
    tbl_height = Inches(0.4) * n_rows
    table_shape = slide.shapes.add_table(n_rows, n_cols, left, top,
                                          sum(col_widths), tbl_height)
    table = table_shape.table
    for i, w in enumerate(col_widths):
        table.columns[i].width = w
    # 表头
    for i, h in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = h
        cell.fill.solid()
        cell.fill.fore_color.rgb = C_TABLE_H
        for p in cell.text_frame.paragraphs:
            p.font.size = Pt(font_size)
            p.font.color.rgb = C_WHITE
            p.font.bold = True
            p.font.name = "微软雅黑"
            p.alignment = PP_ALIGN.CENTER
    # 数据行
    for r_i, row in enumerate(rows):
        for c_i, val in enumerate(row):
            cell = table.cell(r_i + 1, c_i)
            cell.text = str(val)
            cell.fill.solid()
            cell.fill.fore_color.rgb = C_TABLE_R if r_i % 2 == 0 else C_DARK
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(font_size)
                p.font.color.rgb = C_LIGHT
                p.font.name = "微软雅黑"
                p.alignment = PP_ALIGN.CENTER
    # 细边框 (via XML)
    from lxml import etree
    ns_a = 'http://schemas.openxmlformats.org/drawingml/2006/main'
    nsmap = {'a': ns_a}
    for r in range(n_rows):
        for c in range(n_cols):
            cell = table.cell(r, c)
            tc = cell._tc
            tcPr = tc.get_or_add_tcPr()
            for ln_tag in ['lnL', 'lnR', 'lnT', 'lnB']:
                ln = tcPr.find(f'a:{ln_tag}', nsmap)
                if ln is None:
                    ln = etree.SubElement(tcPr, etree.QName(ns_a, ln_tag))
                ln.set('w', '6350')  # 0.5pt in EMU
                solidFill = ln.find('a:solidFill', nsmap)
                if solidFill is None:
                    solidFill = etree.SubElement(ln, etree.QName(ns_a, 'solidFill'))
                srgbClr = solidFill.find('a:srgbClr', nsmap)
                if srgbClr is None:
                    srgbClr = etree.SubElement(solidFill, etree.QName(ns_a, 'srgbClr'))
                srgbClr.set('val', '3A4A6A')

def add_arch_layer(slide, left, top, width, height, label, color, font_size=13):
    """架构分层块"""
    shape = add_rect(slide, left, top, width, height, color)
    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = label
    p.font.size = Pt(font_size)
    p.font.color.rgb = C_WHITE
    p.font.bold = True
    p.font.name = "微软雅黑"
    p.alignment = PP_ALIGN.CENTER
    return shape

def add_icon_card(slide, left, top, width, height, title, desc, color=C_PRIMARY):
    """图标卡片"""
    shape = add_rect(slide, left, top, width, height, C_DARK2)
    # 顶部色条
    add_rect(slide, left, top, width, Pt(3), color)
    add_text_box(slide, left + Inches(0.15), top + Inches(0.15), width - Inches(0.3), Inches(0.4),
                 title, font_size=15, color=color, bold=True)
    add_text_box(slide, left + Inches(0.15), top + Inches(0.55), width - Inches(0.3), height - Inches(0.7),
                 desc, font_size=11, color=C_LIGHT)

TOTAL_SLIDES = 22

# ================================================================
# Slide 1: 封面
# ================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
add_bg(slide, C_DARK)
add_rect(slide, Inches(0), Inches(0), W, Inches(0.08), C_PRIMARY)
add_rect(slide, Inches(0), Inches(7.42), W, Inches(0.08), C_ACCENT)

add_text_box(slide, Inches(1), Inches(1.5), Inches(11), Inches(1.2),
             "SM2 安全数据交换 SDK", font_size=48, color=C_WHITE, bold=True,
             alignment=PP_ALIGN.CENTER)

add_text_box(slide, Inches(1), Inches(2.7), Inches(11), Inches(0.7),
             "基于国密算法的对等系统间安全通信解决方案",
             font_size=24, color=C_ACCENT, alignment=PP_ALIGN.CENTER)

# 分隔线
add_rect(slide, Inches(4), Inches(3.6), Inches(5.3), Pt(2), C_PRIMARY)

add_text_box(slide, Inches(1), Inches(4.0), Inches(11), Inches(0.5),
             "技术分享", font_size=20, color=C_GRAY, alignment=PP_ALIGN.CENTER)

# 特性标签
tags = ["SM2 密钥协商", "SM4-GCM 加解密", "会话管理", "防重放", "Spring Boot 自动配置"]
x_start = Inches(1.8)
for i, tag in enumerate(tags):
    x = x_start + i * Inches(2.1)
    shape = add_rect(slide, x, Inches(5.0), Inches(1.9), Inches(0.45), C_DARK2, C_BORDER)
    tf = shape.text_frame
    p = tf.paragraphs[0]
    p.text = tag
    p.font.size = Pt(11)
    p.font.color.rgb = C_ACCENT
    p.font.name = "微软雅黑"
    p.alignment = PP_ALIGN.CENTER

add_text_box(slide, Inches(1), Inches(6.5), Inches(11), Inches(0.5),
             "2026年7月", font_size=16, color=C_GRAY, alignment=PP_ALIGN.CENTER)

# ================================================================
# Slide 2: 目录
# ================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, C_DARK)
add_section_header(slide, 0, "目  录")
add_bottom_line(slide)
add_page_number(slide, 2, TOTAL_SLIDES)

agenda = [
    ("01", "背景与目标", "为什么需要这套安全数据交换方案"),
    ("02", "总体架构设计", "对等系统架构 · 分层设计 · 模块依赖"),
    ("03", "HTTP 协议规范", "握手协议 · 业务请求格式 · 防重放 · 密钥失效处理"),
    ("04", "核心模块实现", "会话管理 · SM2 密钥协商 · SM4 加解密 · 网络通信 · 访问控制"),
    ("05", "SDK 使用指南", "Spring Boot 自动配置 · Sm2HttpClient · @Sm2Api 注解"),
    ("06", "工程实践", "密钥管理 · 超时重试 · 熔断降级 · 监控告警"),
    ("07", "版本规划与总结", "演进路线 · 安全加固 · Q&A"),
]

for i, (num, title, desc) in enumerate(agenda):
    y = Inches(1.5) + i * Inches(0.82)
    # 编号
    add_text_box(slide, Inches(0.8), y, Inches(0.6), Inches(0.5),
                 num, font_size=28, color=C_PRIMARY, bold=True)
    add_text_box(slide, Inches(1.5), y + Inches(0.0), Inches(3), Inches(0.4),
                 title, font_size=18, color=C_WHITE, bold=True)
    add_text_box(slide, Inches(1.5), y + Inches(0.35), Inches(10), Inches(0.35),
                 desc, font_size=13, color=C_GRAY)

# ================================================================
# Slide 3: 背景与目标
# ================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, C_DARK)
add_section_header(slide, 1, "背景与目标")
add_bottom_line(slide)
add_page_number(slide, 3, TOTAL_SLIDES)

# 左侧 - 痛点
add_text_box(slide, Inches(0.6), Inches(1.4), Inches(5.5), Inches(0.5),
             "🔴 痛点", font_size=22, color=C_RED, bold=True)

pain_points = [
    "系统间数据交互缺乏统一的安全标准，各自实现加解密逻辑",
    "对接过程中密钥管理混乱，私钥硬编码、明文传输时有发生",
    "对接联调成本高：双方需深入理解国密算法细节才能对接",
    "缺乏会话管理机制，密钥长期不变，存在安全隐患",
]
for i, p in enumerate(pain_points):
    add_text_box(slide, Inches(0.8), Inches(2.0) + i * Inches(0.6), Inches(5.5), Inches(0.5),
                 f"• {p}", font_size=14, color=C_LIGHT)

# 右侧 - 目标
add_text_box(slide, Inches(7.0), Inches(1.4), Inches(5.5), Inches(0.5),
             "🟢 目标", font_size=22, color=C_ACCENT, bold=True)

goals = [
    "提供标准化 SDK，双方只需引入依赖 + 配置即可安全通信",
    "基于国密算法（SM2/SM3/SM4），满足合规要求",
    "对等架构：一个 SDK，同时支持主动调用和被动响应",
    "透明加解密：业务代码只处理明文，SDK 自动完成加解密、握手、续期",
]
for i, g in enumerate(goals):
    add_text_box(slide, Inches(7.2), Inches(2.0) + i * Inches(0.6), Inches(5.5), Inches(0.5),
                 f"• {g}", font_size=14, color=C_LIGHT)

# 中间分隔
add_rect(slide, Inches(6.5), Inches(1.6), Pt(1.5), Inches(4.0), C_BORDER)

# 底部 - 核心理念
add_rect(slide, Inches(0.5), Inches(5.5), Inches(12.3), Inches(1.4), C_DARK2)
add_text_box(slide, Inches(0.8), Inches(5.6), Inches(11), Inches(0.4),
             "核心理念：协议是契约，SDK 是便捷封装", font_size=18, color=C_ACCENT, bold=True)
add_text_box(slide, Inches(0.8), Inches(6.1), Inches(11), Inches(0.6),
             "如果对方不信任你的 SDK，可以照着 HTTP 协议规范自己写实现。SDK 只是这份协议的便捷封装，协议本身才是双方通信的契约。",
             font_size=14, color=C_LIGHT)

# ================================================================
# Slide 4: 总体架构 - 对等系统
# ================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, C_DARK)
add_section_header(slide, 2, "总体架构设计 — 对等系统")
add_bottom_line(slide)
add_page_number(slide, 4, TOTAL_SLIDES)

# 左侧 - 系统A
add_icon_card(slide, Inches(0.5), Inches(1.4), Inches(5.8), Inches(2.3),
              "系统 A", "主动调用：sm2Client.post(\"/api/pay\")\n"
              "被动响应：@Sm2Api → 接收 B 的查询请求",
              C_PRIMARY)

# 右侧 - 系统B
add_icon_card(slide, Inches(7.0), Inches(1.4), Inches(5.8), Inches(2.3),
              "系统 B", "主动调用：sm2Client.get(\"/api/query\")\n"
              "被动响应：@Sm2Api → 接收 A 的支付请求",
              C_ACCENT)

# 中间双向箭头
add_text_box(slide, Inches(3.5), Inches(2.0), Inches(6.3), Inches(0.8),
             "⟷  SM2 握手 + SM4 加密通信  ⟷",
             font_size=20, color=C_WARN, bold=True, alignment=PP_ALIGN.CENTER)

# 设计理念
add_text_box(slide, Inches(0.6), Inches(4.0), Inches(12), Inches(0.5),
             "设计理念：一个 SDK，两种角色", font_size=22, color=C_WHITE, bold=True)

design_ideas = [
    "不在代码层面区分\"客户端 SDK\"和\"服务端 SDK\"",
    "通过 yml 配置决定当前系统开启哪些能力：client.enabled / server.enabled",
    "双方各自引入同一个 JAR，配自己的私钥 + 对方的公钥即可通信",
]
for i, d in enumerate(design_ideas):
    add_text_box(slide, Inches(0.8), Inches(4.6) + i * Inches(0.45), Inches(11), Inches(0.4),
                 f"• {d}", font_size=14, color=C_LIGHT)

# 配置示例
code = """# 系统 A 配置（既做客户端又做服务端）
sm2-sdk:
  my-identity: SYSTEM_A
  my-private-key: /secure/a_private.pem
  peer.SYSTEM_B:
    public-key: /secure/b_public.pem
    server-url: https://api.system-b.com
  client.enabled: true      # 开启主动调用
  server.enabled: true      # 开启被动响应"""
add_code_block(slide, Inches(1.0), Inches(5.8), Inches(11.3), Inches(1.5), code, font_size=10)

# ================================================================
# Slide 5: 分层架构
# ================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, C_DARK)
add_section_header(slide, 2, "总体架构设计 — 分层设计")
add_bottom_line(slide)
add_page_number(slide, 5, TOTAL_SLIDES)

layers = [
    ("业务应用层", "主动调用：sm2Client.get()/post()/put()/delete()\n被动响应：@Sm2Api 注解标记 Controller", C_PRIMARY, 2.5),
    ("SDK 门面层", "Sm2HttpClient（主动调用）+ Sm2ServerFilter（被动响应）", RGBColor(0x3A, 0x6B, 0xE7), 2.0),
    ("会话管理模块（核心）", "会话创建/接收 · 密钥缓存(Caffeine+Redis) · 生命周期管理 · 并发控制", RGBColor(0x4A, 0x7B, 0xF7), 2.0),
    ("SM2 密钥协商", "临时密钥对生成 · 共享密钥计算(ECDH) · KDF 密钥派生(SM3)", RGBColor(0x5A, 0x8B, 0xFF), 2.0),
    ("SM4 加解密", "业务数据加密(GCM) · 认证标签生成/校验 · 会话密钥自动绑定", C_ACCENT, 2.0),
]

# 左侧架构图
for i, (name, desc, color, h_inch) in enumerate(layers):
    y = Inches(1.4) + i * Inches(1.15)
    w = Inches(7.5)
    shape = add_rect(slide, Inches(0.4), y, w, Inches(1.05), color)
    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = Pt(12)
    p = tf.paragraphs[0]
    p.text = name
    p.font.size = Pt(15)
    p.font.color.rgb = C_WHITE
    p.font.bold = True
    p.font.name = "微软雅黑"
    p2 = tf.add_paragraph()
    p2.text = desc
    p2.font.size = Pt(10)
    p2.font.color.rgb = RGBColor(0xDD, 0xDD, 0xDD)
    p2.font.name = "微软雅黑"

# 右侧要点
add_text_box(slide, Inches(8.3), Inches(1.4), Inches(4.5), Inches(0.5),
             "关键依赖关系", font_size=20, color=C_WHITE, bold=True)

key_points = [
    "业务层只依赖 Sm2HttpClient + @Sm2Api",
    "会话管理层依赖密钥协商模块和加解密模块",
    "加解密基于 Hutool 实现 SM2/SM3/SM4",
    "主动调用使用 Hutool HttpUtil，被动响应使用 Servlet Filter",
    "注解驱动 vs 配置兜底：正常用注解，第三方接口用 path.include/exclude",
]
for i, kp in enumerate(key_points):
    add_text_box(slide, Inches(8.3), Inches(2.0) + i * Inches(0.55), Inches(4.5), Inches(0.5),
                 f"▸ {kp}", font_size=13, color=C_LIGHT)

# ================================================================
# Slide 6: 协议规范概览
# ================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, C_DARK)
add_section_header(slide, 3, "HTTP 协议规范 — 整体约定")
add_bottom_line(slide)
add_page_number(slide, 6, TOTAL_SLIDES)

add_table(slide, Inches(0.5), Inches(1.4),
          [Inches(2.8), Inches(9.5)],
          ["项目", "约定"],
          [
              ("传输协议", "HTTPS（必须）"),
              ("算法", "SM2 密钥协商 (GB/T 32918.3) + SM4-GCM 加解密"),
              ("编码", "所有二进制字段使用 Base64 编码传输"),
              ("字符编码", "UTF-8"),
              ("时间戳", "Unix 毫秒时间戳"),
              ("SM4 模式", "GCM（认证加密），IV 12字节 / TAG 16字节"),
              ("AAD", "sessionId 的 UTF-8 字节（不随请求体传输，从 Header 提取）"),
              ("防重放", "握手层：时间戳偏差 ≤ 300s；业务层：Nonce 5分钟去重"),
              ("协议版本", "握手 protocolVersion 字段协商，不兼容降级重试"),
          ], font_size=12)

add_text_box(slide, Inches(0.6), Inches(5.7), Inches(12), Inches(0.5),
             "通用请求格式（所有 HTTP 方法统一）", font_size=18, color=C_ACCENT, bold=True)

code_req = """{GET|POST|PUT|DELETE} /api/xxx    ← HTTP 方法和路径与明文接口保持一致
Content-Type: text/plain
X-Session-Id: {sessionId}
X-Timestamp: {unixMillis}
X-Nonce: {16字节随机数 Base64}

Body: BASE64( IV[12字节] || CIPHERTEXT || TAG[16字节] )"""
add_code_block(slide, Inches(0.5), Inches(6.1), Inches(12.3), Inches(1.2), code_req, font_size=11)

# ================================================================
# Slide 7: 握手协议
# ================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, C_DARK)
add_section_header(slide, 3, "HTTP 协议规范 — 三次握手")
add_bottom_line(slide)
add_page_number(slide, 7, TOTAL_SLIDES)

# 三步流程
steps = [
    ("第一步：客户端 → 服务端", "POST /sm2/handshake/init",
     "发送：ephemeralPublicKey(RA) + clientId + ZA + timestamp + signature\n"
     "签名内容：dA 对 (RA || clientId || ZA || timestamp) 签名"),
    ("第二步：服务端 → 客户端", "HTTP 200",
     "返回：sessionId + ephemeralPublicKey(RB) + confirmation(SB)\n"
     "已计算共享密钥并派生 SM4 会话密钥"),
    ("第三步：客户端 → 服务端", "POST /sm2/handshake/confirm",
     "发送：sessionId + confirmation(SA)\n"
     "服务端验证 SA → 握手成功 → 后续用 sessionId 发业务请求"),
]

for i, (title, endpoint, desc) in enumerate(steps):
    y = Inches(1.4) + i * Inches(1.8)
    add_text_box(slide, Inches(0.6), y, Inches(12), Inches(0.4),
                 title, font_size=18, color=C_PRIMARY, bold=True)
    add_text_box(slide, Inches(0.6), y + Inches(0.4), Inches(12), Inches(0.35),
                 endpoint, font_size=14, color=C_ACCENT)
    add_text_box(slide, Inches(0.6), y + Inches(0.75), Inches(12), Inches(0.9),
                 desc, font_size=13, color=C_LIGHT)

# 底部要点
add_rect(slide, Inches(0.5), Inches(6.6), Inches(12.3), Inches(0.7), C_DARK2)
add_text_box(slide, Inches(0.8), Inches(6.65), Inches(11.5), Inches(0.6),
             "💡 关键：三次握手后，双方各自计算出相同的共享密钥，但密钥本身从未在网络中传输。"
             "中间人无法从截获的公钥反推私钥（椭圆曲线离散对数难题）。",
             font_size=13, color=C_WARN)

# ================================================================
# Slide 8: 业务请求格式速查
# ================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, C_DARK)
add_section_header(slide, 3, "HTTP 协议规范 — 业务请求格式速查")
add_bottom_line(slide)
add_page_number(slide, 8, TOTAL_SLIDES)

add_table(slide, Inches(0.5), Inches(1.4),
          [Inches(1.6), Inches(5.0), Inches(5.7)],
          ["HTTP 方法", "Body 明文格式", "说明"],
          [
              ("GET", '{"key1":"value1","key2":"value2"}', "Query 参数键值对 → JSON，无参数时为空字符串"),
              ("POST", '{"field1":"value1","nested":{"a":1}}', "原始 JSON Body，无 Body 时为空字符串"),
              ("POST(文件)", '{"field1":"value1","fileData":"BASE64..."}', "表单字段 + 文件 Base64 合入一个 JSON"),
              ("PUT", '{"field1":"newValue1",...}', "原始 JSON Body，与 POST 规则一致"),
              ("DELETE", '{"id":"xxx","reason":"..."}', "删除条件 → JSON，无参数时为空字符串"),
          ], font_size=12)

add_text_box(slide, Inches(0.6), Inches(4.2), Inches(12), Inches(0.5),
             "统一原则", font_size=18, color=C_ACCENT, bold=True)

add_text_box(slide, Inches(0.8), Inches(4.7), Inches(11.5), Inches(0.7),
             "无论哪种 HTTP 方法，Body 的明文始终是合法的 JSON 字符串（或空字符串），"
             "加密后按 BASE64(IV || CIPHERTEXT || TAG) 格式放入 Body，Content-Type 统一为 text/plain。",
             font_size=14, color=C_LIGHT)

# 请求示例
code_example = """// GET 请求（查询参数加密在 Body）
GET /api/user/query
Content-Type: text/plain
X-Session-Id: a1b2c3d4-e5f6-...
X-Timestamp: 1719676800000
X-Nonce: RfG3kL9mP2sV7wX1Y5zA==
Body: s8dF2kL9mP3xR7wX1yA5zB6cD0eF4gH8...（Base64 密文）"""
add_code_block(slide, Inches(0.5), Inches(5.4), Inches(7.0), Inches(1.6), code_example, font_size=10)

# 响应格式
add_text_box(slide, Inches(8.0), Inches(4.2), Inches(5), Inches(0.5),
             "业务响应格式", font_size=18, color=C_ACCENT, bold=True)

code_resp = """HTTP 200
Content-Type: text/plain
X-Session-Expired: false

Body: BASE64( IV[12B] ||
       CIPHERTEXT ||
       TAG[16B] )"""
add_code_block(slide, Inches(8.0), Inches(4.7), Inches(5.0), Inches(2.3), code_resp, font_size=10)

# ================================================================
# Slide 9: 密钥失效处理
# ================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, C_DARK)
add_section_header(slide, 3, "HTTP 协议规范 — 密钥失效处理")
add_bottom_line(slide)
add_page_number(slide, 9, TOTAL_SLIDES)

add_text_box(slide, Inches(0.6), Inches(1.4), Inches(12), Inches(0.5),
             "三种失效场景 + 自动恢复（对调用方完全透明）",
             font_size=20, color=C_ACCENT, bold=True)

# 场景卡片
scenarios = [
    ("场景一：会话自然过期", C_RED,
     "服务端返回 HTTP 401 + X-Session-Expired: true\n"
     "→ 客户端自动重新握手（3 次 HTTPS）\n"
     "→ 用新会话重试原请求\n"
     "→ 总耗时约 50ms + 原请求耗时"),
    ("场景二：竞争窗口过期", C_WARN,
     "客户端本地缓存显示会话还有 2 分钟\n"
     "请求到达服务端时刚好过期\n"
     "→ 处理方式同场景一，自动恢复"),
    ("场景三：TAG 校验失败", C_RED,
     "密钥不一致或被篡改\n"
     "返回 400 + errorCode=21202\n"
     "→ 废弃当前会话 → 重新握手 → 重试\n"
     "（不能简单重试——密钥仍会失败）"),
]

for i, (title, color, desc) in enumerate(scenarios):
    x = Inches(0.5) + i * Inches(4.2)
    add_icon_card(slide, x, Inches(2.1), Inches(3.9), Inches(2.3), title, desc, color)

# 重试策略表
add_text_box(slide, Inches(0.6), Inches(4.7), Inches(12), Inches(0.5),
             "客户端自动重试协议约定", font_size=18, color=C_WHITE, bold=True)

add_table(slide, Inches(0.5), Inches(5.2),
          [Inches(3.6), Inches(3.5), Inches(2.8), Inches(2.4)],
          ["服务端返回", "含义", "客户端行为", "重试原请求"],
          [
              ("HTTP 401 + X-Session-Expired: true", "会话过期/不存在", "重握手 → 重试", "✅"),
              ("HTTP 400 + errorCode=21202(TAG失败)", "密钥不一致", "废弃会话 → 重握手 → 重试", "✅"),
              ("HTTP 400 + errorCode=21103(签名失败)", "握手阶段证书问题", "检查本方密钥配置", "❌"),
              ("HTTP 403 + Nonce 重复", "重放攻击", "生成新 Nonce 重试一次", "✅(一次)"),
          ], font_size=11)

# ================================================================
# Slide 10: 会话管理模块
# ================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, C_DARK)
add_section_header(slide, 4, "核心模块 — 会话管理")
add_bottom_line(slide)
add_page_number(slide, 10, TOTAL_SLIDES)

# 左侧要点
add_text_box(slide, Inches(0.6), Inches(1.4), Inches(6.5), Inches(0.5),
             "会话生命周期管理", font_size=20, color=C_WHITE, bold=True)

session_features = [
    "本地缓存：ConcurrentHashMap + 过期时间戳（客户端）",
    "分布式缓存：Redis（服务端多实例共享）",
    "惰性删除：获取时检查过期 → 自动删除",
    "定时清理：后台线程每 60s 扫描清理超时会话",
    "自动续期：剩余有效期 < 5min 时自动续期",
    "Rekey：续期时派生新 SM4 密钥，防止密钥过度使用",
    "总生命周期上限：24h，超限强制重握手",
    "限次策略：单会话最多 1000 次请求",
    "并发控制：同一 sessionId 加锁防止重复握手",
]
for i, f in enumerate(session_features):
    add_text_box(slide, Inches(0.6), Inches(2.0) + i * Inches(0.42), Inches(6.5), Inches(0.4),
                 f"▸ {f}", font_size=13, color=C_LIGHT)

# 右侧 - 缓存架构
add_text_box(slide, Inches(7.5), Inches(1.4), Inches(5.5), Inches(0.5),
             "双级缓存架构", font_size=20, color=C_WHITE, bold=True)

add_icon_card(slide, Inches(7.5), Inches(2.0), Inches(5.3), Inches(1.6),
              "本地 Caffeine 缓存", "客户端主动调用时使用\n存储自己的会话密钥\n低延迟，进程内访问", C_PRIMARY)
add_icon_card(slide, Inches(7.5), Inches(3.9), Inches(5.3), Inches(1.6),
              "Redis 分布式缓存", "服务端被动响应时使用\n多实例共享会话\n键: sm2:session:{sessionId}\nSM4 加密存储，密钥仅存内存", C_ACCENT)

add_icon_card(slide, Inches(7.5), Inches(5.8), Inches(5.3), Inches(1.2),
              "Redis 安全存储", "Redis 中的 Session 对象用本地随机生成的\nSM4 密钥加密后存储\n密钥仅存在于内存，重启后旧会话自动失效",
              C_WARN)

# ================================================================
# Slide 11: SM2 密钥协商
# ================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, C_DARK)
add_section_header(slide, 4, "核心模块 — SM2 密钥协商")
add_bottom_line(slide)
add_page_number(slide, 11, TOTAL_SLIDES)

add_text_box(slide, Inches(0.6), Inches(1.4), Inches(12), Inches(0.5),
             "基于 GB/T 32918.3-2016 标准密钥交换协议，实现 ECDH 密钥协商",
             font_size=16, color=C_GRAY)

# ECDH 原理
add_text_box(slide, Inches(0.6), Inches(1.9), Inches(6), Inches(0.4),
             "ECDH 核心原理（颜色混合类比）", font_size=18, color=C_PRIMARY, bold=True)

code_ecdh = """公开基点 G = 黄色颜料（所有人相同）
Alice 私钥 dA = 红色  →  公钥 PA = 红+黄 = 橙色（可公开）
Bob   私钥 dB = 蓝色  →  公钥 PB = 蓝+黄 = 绿色（可公开）

Alice 计算: dA × PB = 红 + 绿 = 红+蓝+黄 = 棕色  }
Bob   计算: dB × PA = 蓝 + 橙 = 蓝+红+黄 = 棕色  } 相同！

数学本质：dA × (dB × G) = dB × (dA × G)
运算顺序可交换，但不可逆（椭圆曲线离散对数难题）"""
add_code_block(slide, Inches(0.5), Inches(2.4), Inches(7.0), Inches(2.8), code_ecdh, font_size=11)

# KDF 密钥派生
add_text_box(slide, Inches(8.0), Inches(1.9), Inches(5), Inches(0.4),
             "KDF 密钥派生 (GB/T 32918.4)", font_size=18, color=C_ACCENT, bold=True)

code_kdf = """输入：Z = x1 || ZA || ZB
      klen = 480 位
算法：
  ct = 0x00000001
  对 i = 1 to ceil(klen/256):
    Ha[i] = SM3(Z || ct)
    ct++

派生输出：
  偏移 0    → SM4 会话密钥 (128位)
  偏移 128  → SM4 IV 初始值 (96位)
  偏移 224  → HMAC-SM3 密钥 (256位, 可选)"""
add_code_block(slide, Inches(8.0), Inches(2.4), Inches(5.0), Inches(4.1), code_kdf, font_size=10)

# 关键参数表
add_table(slide, Inches(0.5), Inches(5.5),
          [Inches(2.0), Inches(4.0), Inches(6.3)],
          ["参数", "默认值", "说明"],
          [
              ("w (椭圆曲线)", "127", "SM2 推荐曲线 w = ⌈log₂(n)⌉ - 1"),
              ("时间戳容差", "300s", "超过偏差则拒绝握手"),
              ("可辨别标识", "双方自定义", "IDA/IDB 用于计算 ZA/ZB"),
              ("派生总长度 klen", "480位", "16+12+32 字节"),
          ], font_size=11)

# ================================================================
# Slide 12: SM4 加解密 + 网络通信
# ================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, C_DARK)
add_section_header(slide, 4, "核心模块 — SM4 加解密 & 网络通信")
add_bottom_line(slide)
add_page_number(slide, 12, TOTAL_SLIDES)

# SM4 模块 - 左
add_text_box(slide, Inches(0.6), Inches(1.4), Inches(6.0), Inches(0.5),
             "SM4-GCM 加解密模块", font_size=20, color=C_PRIMARY, bold=True)

sm4_pts = [
    "强制使用 GCM 模式（认证加密，自带完整性校验）",
    "绝对禁止 ECB 模式；环境不支持 GCM 则降级 CBC+HMAC-SM3",
    "密钥 128位 + GCM IV 96位 + TAG 128位",
    "AAD 固定为 sessionId.getBytes(UTF_8)，从 Header 提取",
    "传输格式：Base64(iv + ciphertext + tag)",
    "解密时自动 TAG 校验，失败抛 IntegrityException",
]
for i, p in enumerate(sm4_pts):
    add_text_box(slide, Inches(0.6), Inches(2.0) + i * Inches(0.42), Inches(6.5), Inches(0.4),
                 f"▸ {p}", font_size=13, color=C_LIGHT)

# 接口设计表
add_text_box(slide, Inches(0.6), Inches(4.6), Inches(6.0), Inches(0.4),
             "两层接口设计", font_size=16, color=C_ACCENT, bold=True)

add_table(slide, Inches(0.5), Inches(5.1),
          [Inches(1.8), Inches(2.2), Inches(2.3)],
          ["层级", "输入", "输出"],
          [
              ("底层(纯加密)", "key, iv, aad, plaintext", "ciphertext + tag"),
              ("上层(业务用)", "sessionId, plaintext", "Base64 密文"),
          ], font_size=11)

# 网络通信 - 右
add_text_box(slide, Inches(7.3), Inches(1.4), Inches(5.5), Inches(0.5),
             "网络通信层设计", font_size=20, color=C_ACCENT, bold=True)

add_text_box(slide, Inches(7.3), Inches(2.0), Inches(5.5), Inches(1.0),
             "核心原则：仅指定的接口走加解密通道，不影响项目内其他 HTTP 通信",
             font_size=14, color=C_WARN)

comm_pts = [
    "服务端：@Sm2Api 注解标记 → Filter 自动加解密",
    "客户端：Sm2HttpClient 工具类 → 内部用 Hutool HttpUtil",
    "无注解的接口 → 直接放行，明文透传",
    "项目中 RestTemplate/Feign/OkHttp 不受影响",
]
for i, p in enumerate(comm_pts):
    add_text_box(slide, Inches(7.3), Inches(2.7) + i * Inches(0.42), Inches(5.5), Inches(0.4),
                 f"▸ {p}", font_size=13, color=C_LIGHT)

# 一次调用全链路
add_text_box(slide, Inches(7.3), Inches(4.6), Inches(5.5), Inches(0.4),
             "一次调用全链路", font_size=16, color=C_ACCENT, bold=True)

code_flow = """系统 A                          系统 B
sm2Client.post("/api/pay")
  ├─ 自动握手 ──────────────────→
  ├─ SM4-GCM 加密 body          │
  ├─ Hutool HttpUtil 发送 ──────→
  │                              ├─ @Sm2Api 检测
  │                              ├─ 解密 → Controller
  │                              ├─ 加密 response
  │  ←──────────────────────────
  └─ 解密 response → Result"""
add_code_block(slide, Inches(7.3), Inches(5.1), Inches(5.5), Inches(2.2), code_flow, font_size=9.5)

# ================================================================
# Slide 13: 接口访问控制
# ================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, C_DARK)
add_section_header(slide, 4, "核心模块 — 接口访问控制")
add_bottom_line(slide)
add_page_number(slide, 13, TOTAL_SLIDES)

add_text_box(slide, Inches(0.6), Inches(1.4), Inches(12), Inches(0.5),
             "不同客户端有不同接口访问权限，由服务端统一管控",
             font_size=18, color=C_GRAY)

# SPI 接口设计
add_text_box(slide, Inches(0.6), Inches(2.0), Inches(6), Inches(0.4),
             "SPI 扩展接口", font_size=18, color=C_PRIMARY, bold=True)

code_spi = """public interface ApiAccessProvider {
    boolean check(String clientId, String path, String method);
    List<String> getAllowedPaths(String clientId);
    default void refresh() {}
}"""
add_code_block(slide, Inches(0.5), Inches(2.5), Inches(6.3), Inches(1.5), code_spi, font_size=11)

add_text_box(slide, Inches(0.6), Inches(4.2), Inches(6), Inches(0.4),
             "两种实现方式", font_size=16, color=C_ACCENT, bold=True)

add_icon_card(slide, Inches(0.5), Inches(4.7), Inches(3.0), Inches(2.2),
              "配置文件（默认）", "application.yml 中配置\n支持 Ant 风格通配符\n支持热更新\n零代码即可使用",
              C_PRIMARY)
add_icon_card(slide, Inches(3.8), Inches(4.7), Inches(3.0), Inches(2.2),
              "数据库（自定义）", "实现 ApiAccessProvider\n查询 MySQL/PostgreSQL\nCaffeine 本地缓存\n动态变更实时生效",
              C_ACCENT)

# 配置示例
add_text_box(slide, Inches(7.3), Inches(1.9), Inches(5.5), Inches(0.4),
             "配置文件示例", font_size=18, color=C_PRIMARY, bold=True)

code_access = """sm2:
  access:
    default-policy: deny    # 默认拒绝
    clients:
      client-payment:
        paths:
          - /api/payment/**
          - /api/order/query
      client-report:
        paths:
          - /api/report/**
          - /api/order/query    # 只读
      client-admin:
        paths:
          - /api/**             # 全部"""
add_code_block(slide, Inches(7.3), Inches(2.4), Inches(5.5), Inches(4.5), code_access, font_size=11)

# ================================================================
# Slide 14: SDK 接入概览
# ================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, C_DARK)
add_section_header(slide, 5, "SDK 使用指南 — 接入概览")
add_bottom_line(slide)
add_page_number(slide, 14, TOTAL_SLIDES)

add_text_box(slide, Inches(0.6), Inches(1.4), Inches(12), Inches(0.5),
             "一个依赖 + 一份配置 + 注解/工具类 = 安全通信",
             font_size=22, color=C_ACCENT, bold=True)

# 三步接入
steps2 = [
    ("① 引入依赖", "pom.xml 加一个 sm2-sdk-spring-boot-starter\n"
     "Fat JAR ~3MB，内置 Hutool + Caffeine + Jackson", C_PRIMARY),
    ("② 配置 yml", "本方私钥 + 对方公钥 + 按需开启 client/server\n"
     "支持 配置文件 / 环境变量 / 代码编程 三种配置方式", C_ACCENT),
    ("③ 使用", "主动调用：注入 Sm2HttpClient → get()/post()/put()/delete()\n"
     "被动响应：Controller 上加 @Sm2Api → SDK 自动加解密", C_WARN),
]

for i, (title, desc, color) in enumerate(steps2):
    x = Inches(0.5) + i * Inches(4.2)
    add_icon_card(slide, x, Inches(2.1), Inches(3.9), Inches(1.8), title, desc, color)

# 接入检查清单
add_text_box(slide, Inches(0.6), Inches(4.3), Inches(12), Inches(0.4),
             "接入检查清单（极简 3 步）", font_size=18, color=C_WHITE, bold=True)

add_table(slide, Inches(0.5), Inches(4.8),
          [Inches(1.2), Inches(3.0), Inches(4.2), Inches(4.0)],
          ["步骤", "操作", "服务端", "客户端"],
          [
              ("①", "加依赖", "sm2-sdk-spring-boot-starter", "同一个 JAR"),
              ("②", "配 yml", "私钥 + 客户端公钥 + Redis\n+ sm2.access 访问控制规则", "client-id + 自身私钥\n+ 服务端公钥 + server-url"),
              ("③", "使用", "Controller 上加 @Sm2Api", "注入 Sm2HttpClient\n调 get()/post()/put()/delete()"),
          ], font_size=12)

add_text_box(slide, Inches(0.6), Inches(6.6), Inches(12), Inches(0.5),
             "💡 项目中原有的 RestTemplate / Feign / OkHttp 调用完全不受影响，只有走 Sm2HttpClient 的请求才经过 SM2 通道",
             font_size=14, color=C_WARN)

# ================================================================
# Slide 15: 客户端使用
# ================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, C_DARK)
add_section_header(slide, 5, "SDK 使用指南 — 主动调用（客户端）")
add_bottom_line(slide)
add_page_number(slide, 15, TOTAL_SLIDES)

add_text_box(slide, Inches(0.6), Inches(1.4), Inches(12), Inches(0.5),
             "Sm2HttpClient — 工具类方式（推荐），内部使用 Hutool HttpUtil",
             font_size=16, color=C_GRAY)

# 核心 API
add_text_box(slide, Inches(0.6), Inches(1.9), Inches(6), Inches(0.4),
             "核心 API", font_size=18, color=C_PRIMARY, bold=True)

code_client = """@Autowired
private Sm2HttpClient sm2Client;

// ===== GET 请求 =====
UserInfo user = sm2Client
    .get("/api/user/query")
    .param("idCard", "110101199001011234")
    .execute(UserInfo.class);

// ===== POST 请求 =====
Result r = sm2Client
    .post("/api/payment/submit")
    .body(orderRequest)
    .execute(Result.class);

// ===== PUT 请求 =====
sm2Client.put("/api/user/update")
    .body(updateRequest).execute(Void.class);

// ===== DELETE 请求 =====
sm2Client.delete("/api/cache/clear")
    .param("scope", "all").execute(Void.class);"""
add_code_block(slide, Inches(0.5), Inches(2.3), Inches(6.5), Inches(3.7), code_client, font_size=10)

# SDK 自动完成流程
add_text_box(slide, Inches(7.5), Inches(1.9), Inches(5.5), Inches(0.4),
             "SDK 内部自动完成（每次 execute）", font_size=18, color=C_ACCENT, bold=True)

auto_steps = [
    "① 会话检查 → 无有效会话则自动 SM2 握手",
    "② 生成随机 IV → SM4-GCM 加密 body → Base64",
    "③ 注入 Header: X-Session-Id + Timestamp + Nonce",
    "④ 通过 Hutool HttpUtil 发送 HTTPS 请求",
    "⑤ 校验响应 TAG → SM4-GCM 解密 body → 返回明文",
    "⑥ 会话快过期 → 自动后台续期（对调用方透明）",
]
for i, s in enumerate(auto_steps):
    add_text_box(slide, Inches(7.5), Inches(2.4) + i * Inches(0.42), Inches(5.5), Inches(0.4),
                 s, font_size=13, color=C_LIGHT)

# 纯手动模式
add_text_box(slide, Inches(7.5), Inches(5.2), Inches(5.5), Inches(0.4),
             "纯手动模式（无 Spring）", font_size=16, color=C_WARN, bold=True)

code_manual = """Sm2HttpClient client = Sm2HttpClient.create(
    SdkConfig.builder()
        .clientIdentity("CLIENT_001")
        .serverIdentity("SERVER_PAYMENT")
        .clientPrivateKeyPem("/path/to/key.pem")
        .serverPublicKeyPem("/path/to/pub.pem")
        .serverUrl("https://api.xxx.com")
        .build());

UserInfo user = client.get("/api/user/query")
    .param("id", "123").execute(UserInfo.class);"""
add_code_block(slide, Inches(7.5), Inches(5.5), Inches(5.5), Inches(1.8), code_manual, font_size=9)

# ================================================================
# Slide 16: 服务端使用
# ================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, C_DARK)
add_section_header(slide, 5, "SDK 使用指南 — 被动响应（服务端）")
add_bottom_line(slide)
add_page_number(slide, 16, TOTAL_SLIDES)

add_text_box(slide, Inches(0.6), Inches(1.4), Inches(12), Inches(0.5),
             "@Sm2Api 注解 — 服务端不需要写任何加解密代码",
             font_size=16, color=C_GRAY)

# 注解用法
add_text_box(slide, Inches(0.6), Inches(1.9), Inches(6.3), Inches(0.4),
             "@Sm2Api 用法", font_size=18, color=C_PRIMARY, bold=True)

code_anno = """// 方式1：标记在方法上
@RestController
public class PaymentController {

    @Sm2Api   // ← 仅该方法需要加解密
    @PostMapping("/api/payment/submit")
    public Result submit(@RequestBody PaymentRequest req) {
        // req 已经是解密后的明文，直接使用
        return paymentService.process(req);
    }

    // 不加注解 → SM2 过滤器直接放行
    @GetMapping("/internal/health")
    public String health() { return "OK"; }
}

// 方式2：标记在类上 → 整个 Controller 都需要加解密
@Sm2Api
@RestController
@RequestMapping("/api/sensitive")
public class SensitiveController { ... }"""
add_code_block(slide, Inches(0.5), Inches(2.3), Inches(6.8), Inches(4.8), code_anno, font_size=10)

# Filter 执行流程
add_text_box(slide, Inches(7.8), Inches(1.9), Inches(5), Inches(0.4),
             "Filter 执行流程", font_size=18, color=C_ACCENT, bold=True)

code_filter = """请求进入
  │
  ▼
根据 URL 查找 HandlerMethod
  │
  ├─ 方法或类上有 @Sm2Api？
  │    │
  │    ├─ 是 → 提取 sessionId
  │    │       → 查缓存取密钥
  │    │       → 解密 body
  │    │       → 替换为明文
  │    │       → 放行到 Controller
  │    │       ↑ 响应时自动加密 ← 返回后
  │    │
  │    └─ 否 → 直接放行（明文透传）"""
add_code_block(slide, Inches(7.3), Inches(2.3), Inches(5.5), Inches(4.8), code_filter, font_size=10)

# ================================================================
# Slide 17: 安全特性
# ================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, C_DARK)
add_section_header(slide, 6, "安全特性总览")
add_bottom_line(slide)
add_page_number(slide, 17, TOTAL_SLIDES)

# 三层密钥架构
add_text_box(slide, Inches(0.6), Inches(1.4), Inches(7), Inches(0.5),
             "三层密钥架构 — 逐层缩小爆炸半径", font_size=20, color=C_ACCENT, bold=True)

add_table(slide, Inches(0.5), Inches(2.0),
          [Inches(2.0), Inches(2.5), Inches(3.0), Inches(4.8)],
          ["层级", "密钥", "生命周期", "泄露影响"],
          [
              ("L1 静态密钥", "SM2 公私钥对", "长期（手动更换）", "需重新分发公钥，但历史会话安全"),
              ("L2 临时密钥", "SM2 临时密钥对", "单次握手（秒级）", "仅影响当次握手"),
              ("L3 会话密钥", "SM4 对称密钥", "单次会话（默认5分钟）", "仅泄露该会话内的数据"),
          ], font_size=12)

# 安全防护矩阵
add_text_box(slide, Inches(0.6), Inches(4.1), Inches(7), Inches(0.5),
             "安全防护矩阵", font_size=20, color=C_WHITE, bold=True)

add_table(slide, Inches(0.5), Inches(4.65),
          [Inches(2.6), Inches(4.0), Inches(5.7)],
          ["防护项", "机制", "说明"],
          [
              ("密钥传输安全", "SM2 密钥交换 (ECDH)", "私钥从不出现在网络中"),
              ("数据机密性", "SM4-GCM 对称加密", "所有业务数据密文传输"),
              ("数据完整性", "SM3 + GCM TAG 校验", "防篡改、防伪造"),
              ("防重放攻击", "Nonce + 时间戳双重校验", "每次请求唯一标识，Bloom+Redis两级"),
              ("身份认证", "SM2 签名 + ZA/ZB 确认", "双向身份确认，防中间人"),
              ("前向安全", "ECDH 临时密钥", "每次握手新密钥，历史会话不可追溯"),
              ("会话隔离", "独立 SM4 密钥", "会话间数据完全隔离"),
              ("内存安全", "MemoryCleanUtil", "密钥使用后立即 Arrays.fill 清零"),
          ], font_size=11)

# ================================================================
# Slide 18: 工程实践
# ================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, C_DARK)
add_section_header(slide, 6, "工程实践 — 超时重试 & 熔断降级")
add_bottom_line(slide)
add_page_number(slide, 18, TOTAL_SLIDES)

# 超时与重试
add_text_box(slide, Inches(0.6), Inches(1.4), Inches(6.0), Inches(0.5),
             "超时与重试机制", font_size=20, color=C_PRIMARY, bold=True)

timeout_items = [
    "握手超时：5s，失败最多重试 3 次（1s, 2s, 4s 递增）",
    "会话过期：自动触发重握手 → 重试原请求（透明）",
    "重握手失败：抛出 SessionEstablishException",
    "熔断保护：连续 5 次握手失败 → 熔断 30s",
    "POST 幂等保护：Body JSON 中 _idempotencyKey 字段",
]
for i, item in enumerate(timeout_items):
    add_text_box(slide, Inches(0.6), Inches(2.0) + i * Inches(0.42), Inches(6.5), Inches(0.4),
                 f"▸ {item}", font_size=13, color=C_LIGHT)

# 降级策略
add_text_box(slide, Inches(0.6), Inches(4.3), Inches(12), Inches(0.5),
             "降级与逃生策略", font_size=20, color=C_WARN, bold=True)

add_table(slide, Inches(0.5), Inches(4.9),
          [Inches(3.0), Inches(5.2), Inches(4.1)],
          ["场景", "降级策略", "恢复条件"],
          [
              ("KMS 不可用", "使用本地 DEK 缓存文件（SM2 公钥加密）继续服务", "KMS 恢复后自动切回"),
              ("Redis 不可用", "降级为本地内存缓存（单机模式）", "Redis 恢复后切回分布式"),
              ("握手服务整体异常", "返回 HandshakeFailedException\n可配置 fallback.enabled 使用上次会话密钥", "握手成功率恢复后关闭"),
              ("Nonce 校验 Redis 故障", "仅依赖时间戳校验（降低安全级别），记录告警", "Redis 恢复后恢复双重校验"),
              ("极端逃生（人工审批）", "管理员临时关闭加密（明文传输）\n全量审计 + P1 告警 + 最长 10 分钟", "故障修复后立即关闭"),
          ], font_size=11)

# ================================================================
# Slide 19: 异常码体系
# ================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, C_DARK)
add_section_header(slide, 6, "工程实践 — 异常码体系")
add_bottom_line(slide)
add_page_number(slide, 19, TOTAL_SLIDES)

add_text_box(slide, Inches(0.6), Inches(1.4), Inches(12), Inches(0.5),
             "5 位数字编码：ABC(模块) + DD(具体错误)  →  格式：级别·角色·操类型·序号",
             font_size=16, color=C_GRAY)

add_table(slide, Inches(0.5), Inches(2.0),
          [Inches(1.2), Inches(1.0), Inches(1.5), Inches(4.5), Inches(4.1)],
          ["错误码", "级别", "操作类型", "说明", "处理建议"],
          [
              ("11301", "⚠️警告", "会话", "客户端会话已过期", "自动重握手"),
              ("21202", "🔴严重", "加解密", "SM4 解密 TAG 校验失败", "废弃会话，重握手"),
              ("21103", "🔴严重", "握手", "服务端证书验签失败", "检查服务端公钥配置"),
              ("11108", "⚠️警告", "握手", "握手超时", "检查网络，增加超时时间"),
              ("22301", "🔴严重", "会话", "服务端会话不存在或已过期", "返回 401 让客户端重握手"),
              ("29001", "🔴严重", "安全", "Nonce 重复（防重放拦截）", "可能遭遇重放攻击，告警"),
              ("31401", "💀致命", "配置", "客户端私钥未配置", "检查密钥配置"),
              ("19003", "⚠️警告", "安全", "连续握手失败超限，触发熔断", "暂停服务，等待恢复"),
          ], font_size=10.5)

# HTTP 映射
add_text_box(slide, Inches(0.6), Inches(5.5), Inches(12), Inches(0.5),
             "SDK 异常码 → HTTP 状态码映射", font_size=18, color=C_ACCENT, bold=True)

add_table(slide, Inches(0.5), Inches(6.0),
          [Inches(1.8), Inches(3.5), Inches(7.0)],
          ["HTTP 状态码", "SDK 错误码范围", "说明"],
          [
              ("200", "—", "成功"),
              ("400", "21103, 21106, 21202", "客户端请求错误（签名失败/密钥确认失败/数据篡改）"),
              ("401", "11301, 22301", "会话过期/不存在，需要重握手"),
              ("403", "11104, 22102, 29001", "非法密钥/公钥验证失败/Nonce重复"),
              ("408", "11108", "握手超时"),
              ("429", "12303", "请求频率超限"),
              ("500", "19004, 39000", "服务端内部错误"),
          ], font_size=11)

# ================================================================
# Slide 20: 监控告警
# ================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, C_DARK)
add_section_header(slide, 6, "工程实践 — 监控与告警")
add_bottom_line(slide)
add_page_number(slide, 20, TOTAL_SLIDES)

# 监控指标
add_text_box(slide, Inches(0.6), Inches(1.4), Inches(6.0), Inches(0.5),
             "Prometheus 监控指标", font_size=20, color=C_PRIMARY, bold=True)

metrics = [
    ("sm2_handshake_total", "握手总次数（按成功/失败区分）"),
    ("sm2_handshake_duration", "握手耗时分布（P50/P99）"),
    ("sm2_session_active", "当前活跃会话数"),
    ("sm2_encrypt_requests", "加密请求总数"),
    ("sm2_decrypt_errors", "解密失败次数（TAG/过期分类）"),
    ("sm2_replay_blocks", "防重放拦截次数"),
]
for i, (name, desc) in enumerate(metrics):
    add_text_box(slide, Inches(0.6), Inches(2.0) + i * Inches(0.38), Inches(6.5), Inches(0.35),
                 f"▸ {name}", font_size=13, color=C_LIGHT)
    add_text_box(slide, Inches(3.2), Inches(2.0) + i * Inches(0.38), Inches(4.0), Inches(0.35),
                 desc, font_size=12, color=C_GRAY)

# 告警规则
add_text_box(slide, Inches(7.3), Inches(1.4), Inches(5.5), Inches(0.5),
             "告警规则", font_size=20, color=C_WARN, bold=True)

add_table(slide, Inches(7.3), Inches(2.0),
          [Inches(2.7), Inches(2.3), Inches(2.6)],
          ["告警项", "阈值", "处理建议"],
          [
              ("握手失败率 > 10%", "最近 5min", "检查网络/密钥配置"),
              ("TAG校验失败突增", "1min > 10次", "可能遭遇篡改攻击"),
              ("会话数超过上限", "> 10000", "检查会话泄漏"),
              ("Nonce重复率 > 1%", "最近 5min", "可能遭遇重放攻击"),
          ], font_size=11)

# 日志规范
add_text_box(slide, Inches(7.3), Inches(4.0), Inches(5.5), Inches(0.5),
             "日志规范", font_size=20, color=C_ACCENT, bold=True)

add_icon_card(slide, Inches(7.3), Inches(4.5), Inches(2.6), Inches(1.3),
              "✅ 必须记录", "握手开始/成功/失败\n会话创建/续期/销毁\n加解密异常\n防重放拦截",
              C_ACCENT)

add_icon_card(slide, Inches(10.2), Inches(4.5), Inches(2.6), Inches(1.3),
              "❌ 严禁记录", "SM2 私钥\n共享密钥 K\nSM4 会话密钥\n业务明文数据",
              C_RED)

add_text_box(slide, Inches(7.3), Inches(6.1), Inches(5.5), Inches(0.7),
             '✅ 正确：sessionId=abc123, 协商成功, 耗时=12ms\n'
             '❌ 错误：sessionId=abc123, key=3F7A8B..., iv=9E2D...',
             font_size=12, color=C_LIGHT)

# ================================================================
# Slide 21: 版本规划
# ================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, C_DARK)
add_section_header(slide, 7, "版本演进规划 & 安全加固")
add_bottom_line(slide)
add_page_number(slide, 21, TOTAL_SLIDES)

add_text_box(slide, Inches(0.6), Inches(1.4), Inches(6.0), Inches(0.5),
             "版本演进路线", font_size=20, color=C_PRIMARY, bold=True)

add_table(slide, Inches(0.5), Inches(2.0),
          [Inches(1.0), Inches(5.5), Inches(2.5)],
          ["版本", "核心能力", "预计工作量"],
          [
              ("V1.0", "基础 SM2 密钥协商 + SM4-GCM 加解密 + 内存会话管理", "3 人周"),
              ("V1.1", "Redis 分布式会话管理 + 自动续期 + Nonce 防重放", "2 人周"),
              ("V1.2", "熔断降级 + Metrics 监控 + 配置热更新", "2 人周"),
              ("V2.0", "双证书认证 (GM/T 0024) + 硬件密码机 (HSM) 适配", "4 人周"),
          ], font_size=12)

# 安全加固清单
add_text_box(slide, Inches(7.3), Inches(1.4), Inches(5.5), Inches(0.5),
             "安全加固清单", font_size=20, color=C_ACCENT, bold=True)

add_text_box(slide, Inches(7.3), Inches(2.0), Inches(5.5), Inches(0.35),
             "编译期安全", font_size=16, color=C_WARN, bold=True)
add_text_box(slide, Inches(7.5), Inches(2.4), Inches(5.3), Inches(0.8),
             "▸ ProGuard 代码混淆，降低逆向风险\n"
             "▸ 敏感字符串编译期加密\n"
             "▸ Jar 包 MD5/SHA256 校验",
             font_size=12, color=C_LIGHT)

add_text_box(slide, Inches(7.3), Inches(3.3), Inches(5.5), Inches(0.35),
             "运行期安全", font_size=16, color=C_WARN, bold=True)
add_text_box(slide, Inches(7.5), Inches(3.7), Inches(5.3), Inches(1.2),
             "▸ 临时密钥使用后 Arrays.fill 清零\n"
             "▸ 堆外内存 (ByteBuffer.allocateDirect)\n"
             "▸ -XX:+DisableAttachMechanism 防内存 dump\n"
             "▸ SDK 自身 Jar 包启动时验签防篡改",
             font_size=12, color=C_LIGHT)

add_text_box(slide, Inches(7.3), Inches(5.0), Inches(5.5), Inches(0.35),
             "通信安全", font_size=16, color=C_WARN, bold=True)
add_text_box(slide, Inches(7.5), Inches(5.4), Inches(5.3), Inches(1.0),
             "▸ HTTPS + SM4 双重加密\n"
             "▸ HSTS 强制 HTTPS\n"
             "▸ 证书固定 (Pinning) + 多指纹备份",
             font_size=12, color=C_LIGHT)

# 发布质量门禁
add_text_box(slide, Inches(0.6), Inches(5.2), Inches(6.5), Inches(0.5),
             "发布前质量门禁", font_size=18, color=C_ACCENT, bold=True)

gates = [
    "✅ 单元测试覆盖率 ≥ 80%（核心模块 ≥ 95%）",
    "✅ 集成测试：模拟完整握手+加密流程",
    "✅ 兼容性测试：JDK 8 / 11 / 17",
    "✅ 性能测试：单次握手 ≤ 50ms（不含网络）",
    "✅ 安全扫描：OWASP Dependency Check",
    "✅ 代码规范：Checkstyle / SpotBugs 通过",
]
for i, g in enumerate(gates):
    add_text_box(slide, Inches(0.6), Inches(5.7) + i * Inches(0.3), Inches(6.8), Inches(0.3),
                 g, font_size=12, color=C_LIGHT)

# ================================================================
# Slide 22: 总结与 Q&A
# ================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, C_DARK)
add_rect(slide, Inches(0), Inches(0), W, Inches(0.08), C_PRIMARY)
add_rect(slide, Inches(0), Inches(7.42), W, Inches(0.08), C_ACCENT)

add_text_box(slide, Inches(1), Inches(1.0), Inches(11), Inches(1.0),
             "总结", font_size=42, color=C_WHITE, bold=True, alignment=PP_ALIGN.CENTER)

add_rect(slide, Inches(4), Inches(2.0), Inches(5.3), Pt(2), C_PRIMARY)

summary_items = [
    "协议标准化：HTTP 协议规范独立于 SDK，双方可按协议自己实现",
    "一个 SDK，两种角色：通过配置决定主动调用 / 被动响应能力",
    "透明加解密：业务代码只处理明文，SDK 自动完成握手、加密、续期",
    "安全合规：国密 SM2/SM3/SM4 全链路，满足等保要求",
    "工程完备：会话管理、熔断降级、监控告警、异常码体系一应俱全",
    "极简接入：一个依赖 + 一份配置 + 注解/工具类 = 安全通信",
]

for i, s in enumerate(summary_items):
    add_text_box(slide, Inches(1.5), Inches(2.4) + i * Inches(0.55), Inches(10.3), Inches(0.5),
                 f"✦ {s}", font_size=17, color=C_LIGHT)

add_rect(slide, Inches(4), Inches(5.9), Inches(5.3), Pt(2), C_PRIMARY)

add_text_box(slide, Inches(1), Inches(6.1), Inches(11), Inches(1.0),
             "Q & A", font_size=42, color=C_ACCENT, bold=True, alignment=PP_ALIGN.CENTER)

add_text_box(slide, Inches(1), Inches(7.0), Inches(11), Inches(0.5),
             "感谢聆听  ·  欢迎交流", font_size=18, color=C_GRAY, alignment=PP_ALIGN.CENTER)

# ── 保存 ────────────────────────────────────────────
output_path = os.path.join(os.path.dirname(os.path.abspath(__file__)),
                           "SM2安全数据交换SDK_技术分享.pptx")
prs.save(output_path)
print(f"PPT generated: {output_path}")
print(f"Total slides: {len(prs.slides)}")